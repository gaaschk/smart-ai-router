"""Document text extraction for uploaded files.

Turns a document into plain text that can be injected into any model's context —
so documents work with every model, not just vision ones. Images are NOT handled
here: they're inlined as base64 for vision models at request time.

Supported:
  - PDF (born-digital text layer) via pypdf
  - Word (.docx) via python-docx, PowerPoint (.pptx) via python-pptx,
    Excel (.xlsx) via openpyxl
  - text/* and common code/data types (utf-8 decode)

A scanned/image-only PDF yields little or no text here; that case is handled
later by rasterizing pages to images and routing to a vision model.

Legacy binary Office formats (.doc/.ppt/.xls) are NOT handled — only the modern
OpenXML (zip-based) formats. Callers gate uploads on is_extractable() so an
unreadable file is refused up front rather than silently injected as a
"couldn't read" placeholder.
"""
from __future__ import annotations

import io

# MIME prefixes/types we treat as decodable plain text.
_TEXT_MIME_PREFIXES = ("text/",)
_TEXT_MIME_EXACT = {
    "application/json",
    "application/xml",
    "application/x-yaml",
    "application/yaml",
    "application/javascript",
    "application/x-sh",
    "application/toml",
    "application/x-python",
}

# Modern OpenXML Office MIME types → the extractor that handles each.
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_OFFICE_MIMES = frozenset({_DOCX_MIME, _PPTX_MIME, _XLSX_MIME})


def is_extractable(mime: str) -> bool:
    """True if we can pull text out of this MIME type here (not images)."""
    mime = (mime or "").lower()
    if mime == "application/pdf":
        return True
    if mime in _OFFICE_MIMES:
        return True
    if any(mime.startswith(p) for p in _TEXT_MIME_PREFIXES):
        return True
    return mime in _TEXT_MIME_EXACT


def extract_text(data: bytes, mime: str, *, filename: str = "") -> str:
    """Best-effort text extraction. Returns "" when nothing can be extracted.

    Never raises on malformed input — extraction is best-effort and must not
    break an upload.
    """
    mime = (mime or "").lower()
    if mime == "application/pdf":
        return _extract_pdf(data)
    if mime == _DOCX_MIME:
        return _extract_docx(data)
    if mime == _PPTX_MIME:
        return _extract_pptx(data)
    if mime == _XLSX_MIME:
        return _extract_xlsx(data)
    if is_extractable(mime):
        return _decode_text(data)
    # Unknown type: try a text decode as a last resort (many code files arrive
    # as application/octet-stream); give up quietly if it isn't text.
    if not mime or mime == "application/octet-stream":
        return _decode_text(data)
    return ""


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, ValueError):
            continue
    return ""


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ""
    try:
        reader = PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:  # noqa: BLE001 — one bad page shouldn't fail all
                continue
        return "\n\n".join(p for p in parts if p).strip()
    except Exception:  # noqa: BLE001 — malformed PDF → no text, not a crash
        return ""


def _extract_docx(data: bytes) -> str:
    """Text from a Word .docx: paragraphs plus any table cell text."""
    try:
        import docx  # python-docx
    except ImportError:
        return ""
    try:
        doc = docx.Document(io.BytesIO(data))
    except Exception:  # noqa: BLE001 — malformed docx → no text, not a crash
        return ""
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_pptx(data: bytes) -> str:
    """Text from a PowerPoint .pptx: every shape's text, slide by slide."""
    try:
        from pptx import Presentation  # python-pptx
    except ImportError:
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception:  # noqa: BLE001
        return ""
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            if text.strip():
                lines.append(text.strip())
        if lines:
            slides.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(slides).strip()


def _extract_xlsx(data: bytes) -> str:
    """Text from an Excel .xlsx: each sheet's non-empty rows as tab-joined cells."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception:  # noqa: BLE001
        return ""
    try:
        sheets = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
        return "\n\n".join(sheets).strip()
    finally:
        wb.close()
