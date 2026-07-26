"""Generate document files from markdown-ish text.

The mirror of extract.py: extract.py turns an uploaded document into text so a
model can read it; docgen.py turns a model's text back into a real document the
user can download. It renders one input format — lightweight markdown — into any
of the supported binary/text targets:

  .md / .txt  → written as-is (no rendering)
  .pdf        → reportlab (headings, paragraphs, bullets, tables)
  .docx       → python-docx
  .pptx       → python-pptx (one slide per top-level "# " heading)
  .xlsx       → openpyxl (a markdown table becomes rows; else one line per row)

The markdown we understand is deliberately small — headings (#/##/###), bullet
lines (- or *), pipe tables, blank-line-separated paragraphs, and **bold**
inline. That is enough for the documents an assistant actually produces
(resumes, reports, summaries, simple sheets) without pulling a full markdown
engine.

Rendering is best-effort but, unlike extraction, a failure here DOES matter —
the caller wants a file. Each renderer raises DocGenError on failure so the tool
layer can report a clear message to the model instead of silently producing an
empty file.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass

# Extension → the format key we render. Text formats pass through untouched.
_TEXT_EXTS = {".md", ".txt", ".markdown"}
_DOC_EXTS = {".pdf", ".docx", ".pptx", ".xlsx"}
SUPPORTED_EXTS = _TEXT_EXTS | _DOC_EXTS


class DocGenError(Exception):
    """Rendering a document failed (bad input or a missing optional library)."""


def is_supported(path: str) -> bool:
    return _ext(path) in SUPPORTED_EXTS


def _ext(path: str) -> str:
    dot = path.rfind(".")
    return path[dot:].lower() if dot != -1 else ""


# ── markdown model ────────────────────────────────────────────────────────────

@dataclass
class Block:
    kind: str          # "h1" | "h2" | "h3" | "bullet" | "para" | "table"
    text: str = ""     # for headings/bullets/paragraphs
    rows: list = None  # for tables: list[list[str]]


_HEADING = re.compile(r"^(#{1,3})\s+(.*)$")
_BULLET = re.compile(r"^\s*[-*]\s+(.*)$")
_BOLD = re.compile(r"\*\*(.+?)\*\*")


def _is_table_row(line: str) -> bool:
    return line.strip().startswith("|") or ("|" in line and line.count("|") >= 2)


def _is_table_separator(line: str) -> bool:
    # A markdown header underline like  |---|:--:|---| — every cell is dashes
    # (with optional alignment colons/whitespace), and at least one dash total.
    cells = line.strip().strip("|").split("|")
    if not cells:
        return False
    ok = all(cell.strip() and set(cell.strip()) <= set("-:") for cell in cells)
    return ok and "-" in line


def _split_row(line: str) -> list[str]:
    cells = line.strip().strip("|").split("|")
    return [_strip_inline(c.strip()) for c in cells]


def _strip_inline(text: str) -> str:
    """Drop markdown emphasis markers for plain-text targets."""
    text = _BOLD.sub(r"\1", text)
    return text.replace("`", "")


def parse_markdown(content: str) -> list[Block]:
    """Parse our small markdown subset into an ordered list of blocks.

    Inline markup (**bold**) is preserved in the block text; renderers that can
    show emphasis (PDF, docx) convert it, and plain-text renderers strip it via
    _strip_inline. Table cells are stripped here since no target renders emphasis
    inside a cell.
    """
    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    blocks: list[Block] = []
    para: list[str] = []
    table: list[list[str]] = []

    def flush_para():
        if para:
            blocks.append(Block("para", " ".join(para).strip()))
            para.clear()

    def flush_table():
        if table:
            blocks.append(Block("table", rows=[r for r in table]))
            table.clear()

    for line in lines:
        if _is_table_row(line):
            if _is_table_separator(line):
                continue  # header/body divider — no data
            flush_para()
            table.append(_split_row(line))
            continue
        flush_table()

        if not line.strip():
            flush_para()
            continue
        m = _HEADING.match(line)
        if m:
            flush_para()
            level = len(m.group(1))
            blocks.append(Block(f"h{level}", m.group(2).strip()))
            continue
        m = _BULLET.match(line)
        if m:
            flush_para()
            blocks.append(Block("bullet", m.group(1).strip()))
            continue
        para.append(line.strip())

    flush_para()
    flush_table()
    return blocks


# ── public entry point ──────────────────────────────────────────────────────

def render(path: str, content: str) -> bytes:
    """Render `content` (markdown) into bytes for the file type implied by
    `path`'s extension. Raises DocGenError for unsupported types or on failure.
    """
    ext = _ext(path)
    if ext in _TEXT_EXTS:
        return content.encode("utf-8")
    if ext == ".pdf":
        return _render_pdf(content)
    if ext == ".docx":
        return _render_docx(content)
    if ext == ".pptx":
        return _render_pptx(content)
    if ext == ".xlsx":
        return _render_xlsx(content)
    raise DocGenError(
        f"unsupported document type {ext!r}; supported: "
        + ", ".join(sorted(SUPPORTED_EXTS))
    )


def _render_pdf(content: str) -> bytes:
    try:
        from reportlab.lib.enums import TA_LEFT
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import (
            ListFlowable,
            ListItem,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )
        from reportlab.lib import colors
    except ImportError as exc:  # pragma: no cover - dep is declared
        raise DocGenError(f"PDF support unavailable: {exc}")

    def _inline(text: str) -> str:
        # reportlab Paragraph accepts a mini-HTML subset — map **bold**.
        escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        return _BOLD.sub(r"<b>\1</b>", escaped)

    styles = getSampleStyleSheet()
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, title="")
    flow = []
    pending_bullets: list = []

    def flush_bullets():
        if pending_bullets:
            flow.append(ListFlowable(list(pending_bullets), bulletType="bullet"))
            flow.append(Spacer(1, 6))
            pending_bullets.clear()

    try:
        for b in parse_markdown(content):
            if b.kind != "bullet":
                flush_bullets()
            if b.kind == "h1":
                flow.append(Paragraph(_inline(b.text), styles["Title"]))
            elif b.kind == "h2":
                flow.append(Paragraph(_inline(b.text), styles["Heading2"]))
            elif b.kind == "h3":
                flow.append(Paragraph(_inline(b.text), styles["Heading3"]))
            elif b.kind == "bullet":
                pending_bullets.append(
                    ListItem(Paragraph(_inline(b.text), styles["BodyText"]))
                )
            elif b.kind == "table":
                t = Table(b.rows)
                t.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                flow.append(t)
                flow.append(Spacer(1, 6))
            else:  # para
                flow.append(Paragraph(_inline(b.text), styles["BodyText"]))
                flow.append(Spacer(1, 6))
        flush_bullets()
        if not flow:
            flow.append(Paragraph("", styles["BodyText"]))
        doc.build(flow)
    except Exception as exc:  # noqa: BLE001 — surface as a clean gen error
        raise DocGenError(f"failed to build PDF: {exc}")
    return buf.getvalue()


def _render_docx(content: str) -> bytes:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - dep is declared
        raise DocGenError(f"Word (.docx) support unavailable: {exc}")

    def _add_runs(paragraph, text: str):
        # Split on **bold** spans, adding bold runs for the emphasized parts.
        pos = 0
        for m in _BOLD.finditer(text):
            if m.start() > pos:
                paragraph.add_run(text[pos:m.start()])
            paragraph.add_run(m.group(1)).bold = True
            pos = m.end()
        if pos < len(text):
            paragraph.add_run(text[pos:])

    try:
        doc = docx.Document()
        for b in parse_markdown(content):
            if b.kind == "h1":
                doc.add_heading(_strip_inline(b.text), level=0)
            elif b.kind == "h2":
                doc.add_heading(_strip_inline(b.text), level=1)
            elif b.kind == "h3":
                doc.add_heading(_strip_inline(b.text), level=2)
            elif b.kind == "bullet":
                _add_runs(doc.add_paragraph(style="List Bullet"), b.text)
            elif b.kind == "table":
                cols = max((len(r) for r in b.rows), default=1)
                table = doc.add_table(rows=0, cols=cols)
                table.style = "Table Grid"
                for row in b.rows:
                    cells = table.add_row().cells
                    for i, val in enumerate(row[:cols]):
                        cells[i].text = val
            else:  # para
                _add_runs(doc.add_paragraph(), b.text)
        buf = io.BytesIO()
        doc.save(buf)
    except Exception as exc:  # noqa: BLE001
        raise DocGenError(f"failed to build .docx: {exc}")
    return buf.getvalue()


def _render_pptx(content: str) -> bytes:
    """One slide per top-level '# ' heading; other blocks become body bullets.

    Content before the first heading (if any) starts an untitled slide, so
    heading-less input still yields something.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - dep is declared
        raise DocGenError(f"PowerPoint (.pptx) support unavailable: {exc}")

    try:
        blocks = parse_markdown(content)
        # Group blocks into (title, [body-lines]) slides split on h1/h2.
        slides: list[tuple[str, list[str]]] = []
        title, body = "", []
        for b in blocks:
            if b.kind in ("h1", "h2"):
                if title or body:
                    slides.append((title, body))
                title, body = _strip_inline(b.text), []
            elif b.kind == "table":
                for row in b.rows:
                    body.append(" | ".join(row))
            else:
                body.append(_strip_inline(b.text))
        if title or body:
            slides.append((title, body))
        if not slides:
            slides = [("", [])]

        prs = Presentation()
        blank = prs.slide_layouts[6]
        for stitle, lines in slides:
            slide = prs.slides.add_slide(blank)
            top = Inches(0.4)
            if stitle:
                box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(1))
                tf = box.text_frame
                tf.text = stitle
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True
                top = Inches(1.5)
            if lines:
                box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(5))
                tf = box.text_frame
                tf.word_wrap = True
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
        buf = io.BytesIO()
        prs.save(buf)
    except Exception as exc:  # noqa: BLE001
        raise DocGenError(f"failed to build .pptx: {exc}")
    return buf.getvalue()


def _render_xlsx(content: str) -> bytes:
    """A markdown pipe table becomes rows/columns; otherwise each non-empty
    line is one row (comma-split into cells)."""
    try:
        from openpyxl import Workbook
    except ImportError as exc:  # pragma: no cover - dep is declared
        raise DocGenError(f"Excel (.xlsx) support unavailable: {exc}")

    try:
        blocks = parse_markdown(content)
        rows: list[list[str]] = []
        tables = [b for b in blocks if b.kind == "table"]
        if tables:
            for t in tables:
                rows.extend(t.rows)
        else:
            for line in content.replace("\r\n", "\n").split("\n"):
                if line.strip():
                    rows.append([c.strip() for c in line.split(",")])

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        for row in rows:
            ws.append(row if row else [""])
        buf = io.BytesIO()
        wb.save(buf)
    except Exception as exc:  # noqa: BLE001
        raise DocGenError(f"failed to build .xlsx: {exc}")
    return buf.getvalue()
